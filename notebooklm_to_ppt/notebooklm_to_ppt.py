"""
notebooklm_to_ppt.py
====================
將 NotebookLM 圖片型 PDF 轉為可編輯 PPTX。
不需要任何 API Key，完全離線執行。

流程：
  1. PDF 每頁 → PNG
  2. EasyOCR 提取文字 + 座標
  3. 依座標推斷標題 / 區塊 / 條列結構
  4. 輸出統一字型的 PPTX

使用方式：
  python notebooklm_to_ppt.py 123.pdf
  python notebooklm_to_ppt.py 123.pdf output.pptx
"""

import sys
import os
import numpy as np
from pathlib import Path

import fitz                  # PyMuPDF
import easyocr
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

import config

# ──────────────────────────────────────────────
# 顏色常數
# ──────────────────────────────────────────────
C_TITLE_BG   = RGBColor(31,  73,  125)
C_WHITE      = RGBColor(255, 255, 255)
C_DARK       = RGBColor(40,  40,  40)
C_MID        = RGBColor(110, 110, 110)
C_LIGHT      = RGBColor(240, 242, 248)
C_ACCENT     = RGBColor(31,  73,  125)
C_BORDER     = RGBColor(180, 200, 230)

# ──────────────────────────────────────────────
# Step 1：PDF → PNG bytes list
# ──────────────────────────────────────────────

def pdf_to_images(pdf_path: str) -> list[tuple[bytes, int, int]]:
    """回傳 [(png_bytes, width_px, height_px), ...]"""
    doc = fitz.open(pdf_path)
    result = []
    mat = fitz.Matrix(2.0, 2.0)
    for page in doc:
        pix = page.get_pixmap(matrix=mat)
        result.append((pix.tobytes("png"), pix.width, pix.height))
    doc.close()
    print(f"  共 {len(result)} 頁")
    return result


# ──────────────────────────────────────────────
# Step 2：EasyOCR 提取文字與座標
# ──────────────────────────────────────────────

def init_ocr() -> easyocr.Reader:
    print("  載入 OCR 模型（首次執行會下載，約 1–2 分鐘）...")
    # ch_tra = 繁體中文, en = 英文
    reader = easyocr.Reader(["ch_tra", "en"], gpu=False, verbose=False)
    print("  OCR 模型就緒")
    return reader


def ocr_page(reader: easyocr.Reader, png_bytes: bytes) -> list[dict]:
    """
    回傳 list of {text, cx, cy, height, conf}
    cx/cy 為文字框中心點（0~1 的比例座標）
    height 為文字框高度比例
    """
    results = reader.readtext(png_bytes, detail=1, paragraph=False)
    # results: [ ([[x1,y1],[x2,y2],[x3,y3],[x4,y4]], text, conf), ... ]
    # 我們先取得圖片尺寸（用 fitz 已知，改由傳入）— 這裡從 bbox 相對估算
    if not results:
        return []

    # 取整張圖的高寬上界
    all_xs = [pt[0] for r in results for pt in r[0]]
    all_ys = [pt[1] for r in results for pt in r[0]]
    img_w = max(all_xs) if all_xs else 1
    img_h = max(all_ys) if all_ys else 1

    blocks = []
    for bbox, text, conf in results:
        text = text.strip()
        if not text or conf < 0.3:
            continue
        xs = [p[0] for p in bbox]
        ys = [p[1] for p in bbox]
        cx = (min(xs) + max(xs)) / 2 / img_w
        cy = (min(ys) + max(ys)) / 2 / img_h
        h  = (max(ys) - min(ys)) / img_h
        blocks.append({"text": text, "cx": cx, "cy": cy, "height": h, "conf": conf})

    return blocks


# ──────────────────────────────────────────────
# Step 3：依座標推斷結構
# ──────────────────────────────────────────────

def infer_structure(blocks: list[dict]) -> dict:
    """
    將 OCR 結果依位置推斷為：
      title / subtitle / sections(heading + points) / footer
    """
    if not blocks:
        return {"title": "", "subtitle": "", "sections": [], "footer": ""}

    # 依 cy 排序
    blocks = sorted(blocks, key=lambda b: b["cy"])

    # 找最大字高（通常是標題）
    max_h = max(b["height"] for b in blocks)

    # 分層閾值
    TITLE_ZONE   = 0.25   # cy < 0.25 且字高大 → 標題區
    FOOTER_ZONE  = 0.88   # cy > 0.88 → 頁腳

    title_blocks   = []
    body_blocks    = []
    footer_blocks  = []

    for b in blocks:
        if b["cy"] > FOOTER_ZONE:
            footer_blocks.append(b)
        elif b["cy"] < TITLE_ZONE and b["height"] >= max_h * 0.55:
            title_blocks.append(b)
        else:
            body_blocks.append(b)

    # 標題：合併同列文字
    title = _merge_line(title_blocks)

    # 副標題：緊接在標題下方的次大文字
    subtitle = ""
    if body_blocks:
        sub_candidates = [b for b in body_blocks if b["cy"] < 0.35]
        if sub_candidates:
            subtitle = _merge_line(sub_candidates)
            body_blocks = [b for b in body_blocks if b["cy"] >= 0.35]

    # 頁腳
    footer = " ".join(b["text"] for b in footer_blocks).strip()

    # Body → 依 cx 分欄（偵測欄位數）
    sections = _cluster_into_sections(body_blocks)

    return {
        "title": title,
        "subtitle": subtitle,
        "sections": sections,
        "footer": footer,
    }


def _merge_line(blocks: list[dict]) -> str:
    """將同一區域的文字依 cx 排序後合併"""
    if not blocks:
        return ""
    return " ".join(b["text"] for b in sorted(blocks, key=lambda b: b["cx"]))


def _cluster_into_sections(blocks: list[dict]) -> list[dict]:
    """依 cx（水平位置）分欄，每欄再依 cy 分為 heading + points"""
    if not blocks:
        return []

    # 取得 cx 分布，找欄數
    cxs = np.array([b["cx"] for b in blocks])

    # 簡單分欄：若 cx 全在 0.33~0.66 → 單欄
    # 若有明顯聚集 → 多欄（用固定分位切割）
    cx_range = cxs.max() - cxs.min()

    if cx_range < 0.35:
        # 單欄
        cols = [blocks]
    elif cx_range < 0.65:
        # 雙欄
        mid = (cxs.max() + cxs.min()) / 2
        cols = [
            [b for b in blocks if b["cx"] <= mid],
            [b for b in blocks if b["cx"] >  mid],
        ]
    else:
        # 三欄
        q1 = cxs.min() + cx_range / 3
        q2 = cxs.min() + cx_range * 2 / 3
        cols = [
            [b for b in blocks if b["cx"] <= q1],
            [b for b in blocks if q1 < b["cx"] <= q2],
            [b for b in blocks if b["cx"] >  q2],
        ]

    sections = []
    for col in cols:
        if not col:
            continue
        col_sorted = sorted(col, key=lambda b: b["cy"])

        # 找欄內字高最大的 → heading
        max_h_in_col = max(b["height"] for b in col_sorted)
        heading_blocks = [b for b in col_sorted if b["height"] >= max_h_in_col * 0.75]
        point_blocks   = [b for b in col_sorted if b["height"] <  max_h_in_col * 0.75]

        heading = _merge_line(heading_blocks)
        points  = [b["text"] for b in sorted(point_blocks, key=lambda b: b["cy"])]

        sections.append({"heading": heading, "points": points})

    return sections


# ──────────────────────────────────────────────
# Step 4：建立 PPTX
# ──────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s


def add_text(slide, l, t, w, h, text, size, color,
             bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text          = text
    r.font.name     = config.FONT_NAME
    r.font.size     = Pt(size)
    r.font.color.rgb = color
    r.font.bold     = bold
    r.font.italic   = italic
    return txb


def build_slide(prs, data: dict, num: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W = prs.slide_width
    H = prs.slide_height

    title    = data.get("title", "") or f"第 {num} 頁"
    subtitle = data.get("subtitle", "") or ""
    sections = data.get("sections", []) or []
    footer   = data.get("footer", "") or ""

    # 白底
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_WHITE

    # ── 標題列 ──
    th = Cm(2.6) if not subtitle else Cm(3.5)
    add_rect(slide, Cm(0), Cm(0), W, th, C_TITLE_BG)
    add_text(slide, Cm(0.8), Cm(0.25), W - Cm(2.5), Cm(1.8),
             title, config.FONT_SIZE_TITLE, C_WHITE, bold=True)
    if subtitle:
        add_text(slide, Cm(0.8), Cm(1.9), W - Cm(2.5), Cm(1.2),
                 subtitle, config.FONT_SIZE_SUBTITLE, RGBColor(200, 220, 245))
    # 頁碼
    add_text(slide, W - Cm(2.3), Cm(0.3), Cm(2.0), Cm(0.8),
             f"{num} / {total}", 9, C_WHITE, align=PP_ALIGN.RIGHT)

    # ── 頁腳 ──
    footer_h = Cm(0.85) if footer else Cm(0)
    if footer:
        fy = H - footer_h
        add_rect(slide, Cm(0), fy, W, footer_h, C_LIGHT)
        add_text(slide, Cm(1.0), fy + Cm(0.1), W - Cm(2.0), Cm(0.7),
                 footer, 9, C_MID, italic=True)

    # ── 內容區 ──
    ct = th + Cm(0.4)
    ch = H - ct - footer_h - Cm(0.2)

    if not sections:
        return

    n = len(sections)
    gap = Cm(0.5)

    if n == 1:
        _draw_section(slide, sections[0], Cm(1.0), ct, W - Cm(2.0), ch)
    elif n == 2:
        cw = (W - Cm(2.0) - gap) / 2
        for i, sec in enumerate(sections):
            _draw_section(slide, sec, Cm(1.0) + i * (cw + gap), ct, cw, ch)
    elif n == 3:
        cw = (W - Cm(2.0) - gap * 2) / 3
        for i, sec in enumerate(sections):
            _draw_section(slide, sec, Cm(1.0) + i * (cw + gap), ct, cw, ch)
    else:
        # 4+：雙欄多行
        cw = (W - Cm(2.0) - gap) / 2
        rows = (n + 1) // 2
        rh = (ch - gap * (rows - 1)) / rows
        for i, sec in enumerate(sections):
            col = i % 2
            row = i // 2
            l = Cm(1.0) + col * (cw + gap)
            t = ct + row * (rh + gap)
            _draw_section(slide, sec, l, t, cw, rh)


def _draw_section(slide, sec: dict, l, t, w, h):
    heading = sec.get("heading", "") or ""
    points  = sec.get("points",  []) or []

    # 區塊底色
    add_rect(slide, l, t, w, h, C_LIGHT, C_BORDER)

    cursor = t + Cm(0.2)

    if heading:
        hh = Cm(0.95)
        add_rect(slide, l, t, w, hh, C_ACCENT)
        add_text(slide, l + Cm(0.3), t + Cm(0.1),
                 w - Cm(0.6), hh - Cm(0.1),
                 heading, config.FONT_SIZE_HEADING, C_WHITE, bold=True)
        cursor = t + hh + Cm(0.2)

    if points:
        pts_h = h - (cursor - t) - Cm(0.15)
        txb = slide.shapes.add_textbox(l + Cm(0.3), cursor, w - Cm(0.6), pts_h)
        tf  = txb.text_frame
        tf.word_wrap = True
        for i, pt in enumerate(points):
            pt = pt.strip()
            if not pt:
                continue
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.space_before = Pt(4)
            run = para.add_run()
            run.text          = f"• {pt}"
            run.font.name     = config.FONT_NAME
            run.font.size     = Pt(config.FONT_SIZE_BODY)
            run.font.color.rgb = C_DARK


def build_pptx(all_data: list[dict], output_path: str):
    prs = Presentation()
    prs.slide_width  = Cm(config.SLIDE_WIDTH_CM)
    prs.slide_height = Cm(config.SLIDE_HEIGHT_CM)
    total = len(all_data)
    for i, data in enumerate(all_data, 1):
        build_slide(prs, data, i, total)
    prs.save(output_path)
    print(f"[完成] 已儲存：{output_path}  ({total} 張投影片)")


# ──────────────────────────────────────────────
# 主流程
# ──────────────────────────────────────────────

def main():
    if len(sys.argv) < 2:
        print("用法：python notebooklm_to_ppt.py <input.pdf> [output.pptx]")
        sys.exit(1)

    input_path  = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) >= 3 else \
                  str(Path(input_path).stem) + "_clean.pptx"

    if not os.path.exists(input_path):
        print(f"錯誤：找不到 {input_path}")
        sys.exit(1)

    print(f"[1/3] 讀取 PDF：{input_path}")
    images = pdf_to_images(input_path)

    print(f"[2/3] OCR 文字辨識...")
    reader   = init_ocr()
    all_data = []
    for i, (img_bytes, iw, ih) in enumerate(images, 1):
        print(f"      第 {i}/{len(images)} 頁...", end=" ", flush=True)
        blocks = ocr_page(reader, img_bytes)
        data   = infer_structure(blocks)
        print(f"標題：{data['title'][:30]}  區塊數：{len(data['sections'])}")
        all_data.append(data)

    print(f"[3/3] 建立 PPTX：{output_path}")
    build_pptx(all_data, output_path)


if __name__ == "__main__":
    main()
